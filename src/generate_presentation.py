import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    # Configuration en format 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Charte de couleurs (RGB)
    c_primary = RGBColor(26, 54, 93) # #1A365D - Bleu Profond
    c_secondary = RGBColor(43, 108, 176) # #2B6CB0 - Bleu Secondaire
    c_accent = RGBColor(197, 48, 48) # #C53030 - Rouge Accent
    c_text = RGBColor(45, 55, 72) # #2D3748 - Gris Sombre
    c_muted = RGBColor(113, 128, 150) # #718096 - Gris Neutre

    def add_title(slide, text):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c_primary
        return tb

    def add_bullet_points(slide, points, left, top, width, height, size=16):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, pt_text in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = pt_text
            p.font.name = 'Arial'
            p.font.size = Pt(size)
            p.font.color.rgb = c_text
            p.level = 0
            p.space_after = Pt(10)
        return tb

    def try_add_image(slide, path, left, top, width, height):
        if os.path.exists(path):
            slide.shapes.add_picture(path, left, top, width=width, height=height)
            print(f"📸 Image intégrée sur la diapositive : {path}")
        else:
            # Placeholder visuel en cas d'absence
            tb = slide.shapes.add_textbox(left, top, width, height)
            p = tb.text_frame.paragraphs[0]
            p.text = f"[ Emplacement Graphique :\n{path} non trouvé ]"
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = c_muted
            p.font.size = Pt(12)

    # -------------------------------------------------------------------------
    # DIAPO 1 : TITRE
    # -------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_slide_layout)
    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf1 = tb1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "ÉVALUATION INTÉGRÉE DU POTENTIEL HYDROTHERMAL"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = c_primary
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Caractérisation du Bloc UK Central et Synthèse Réservoir pour le document Rapport_Technique_Reservoir_Geothermie_2.pdf"
    p1_sub.font.size = Pt(18)
    p1_sub.font.color.rgb = c_secondary
    p1_sub.space_before = Pt(15)

    p1_meta = tf1.add_paragraph()
    p1_meta.text = "Direction des Réservoirs & Géologie Numérique | 5 juin 2026"
    p1_meta.font.size = Pt(13)
    p1_meta.font.color.rgb = c_muted
    p1_meta.space_before = Pt(40)

    # -------------------------------------------------------------------------
    # DIAPO 2 : CONTEXTE ET OPPORTUNITÉ STRATÉGIQUE
    # -------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_title(s2, "1. Contexte, Enjeux et Objectifs du Projet")
    pts2_left = [
        "• Stratégie de décarbonation : Ciblage des réseaux de chaleur urbains et des processus industriels au Royaume-Uni.",
        "• Complexité géologique : Superposition d'un prisme sédimentaire permo-triasique syn-to-post-rift sur un socle paléozoïque plissé varisque.",
        "• Objectif de la campagne : S'affranchir des approches empiriques classiques et lever les incertitudes de continuité structurale.",
        "• But final : Sécuriser le doublet hydrothermique à long terme (durée de vie nominale projetée de 30 ans)."
    ]
    add_bullet_points(s2, pts2_left, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), size=17)
    try_add_image(s2, os.path.join("outputs", "plots", "Well_GT01_Tracks.png"), Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.8))

    # -------------------------------------------------------------------------
    # DIAPO 3 : ARCHITECTURE NUMÉRIQUE & DATA PIPELINE
    # -------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_title(s3, "2. Moteur Analytique & Data Pipeline Réel")
    pts3 = [
        "• Automated Core Engine (v4.0.2) : Interconnexion directe et reproductible de l'ingénierie.",
        "• Base de données relationnelle : Extraction en temps réel depuis MySQL Workbench.",
        "• Ingestion de trois piliers fondamentaux de données de sous-sol :",
        " - Données Diagraphiques (Well Logs) : Profils continus GR, Rt, Rhob et NPHI.",
        " - Données Structurales : Épaisseurs brutes (Gross) et géométrie/rejet des accidents majeurs.",
        " - Données Dynamiques : Températures de fond stabilisées (BHT) et pressions initiales."
    ]
    add_bullet_points(s3, pts3, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0), size=16)

    # -------------------------------------------------------------------------
    # DIAPO 4 : ALGORITHMES DE DATA SCIENCE & NET PAY
    # -------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_title(s4, "3. Prétraitement et Logique de Tri des Données")
    pts4 = [
        "• Nettoyage des signaux : Correction des pannes de capteurs et des éboulements de parois par interpolation linéaire bornée.",
        "• Écrêtage automatique des aberrations physiques par le script (Porosité négative, Sw > 100%).",
        "• Application stricte de filtres de coupure (Cut-offs) cumulatifs pour isoler le Net Pay :",
        " - Indice d'argilosité (Gamma Ray) : Vsh < 30%",
        " - Porosité utile (Matrice / Fracture) : Phi > 10%",
        " - Saturation en eau mobile (Milieu hydrothermal) : Sw > 50%",
        "• Résultat : Élimination rigoureuse des intervalles étanches ou non productifs."
    ]
    add_bullet_points(s4, pts4, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0), size=16)

    # -------------------------------------------------------------------------
    # DIAPO 5 : SYNTHÈSE ANALYTIQUE DES PUITS
    # -------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_title(s5, "4. Analyse Comparative de l'Exploration")
    pts5 = [
        "• GT-01 (Bassin Est) : Réservoir matriciel de grès exceptionnel (Net Pay 594 m) mais froid (21.99°C). HIP exploitable nul.",
        "• GT-02 (Axe Central) : Réservoir de fracture restreint (Net Pay 150 m) mais hypertherme (53.74°C). Potentiel commercial majeur.",
        "• GT-03 (Bordure Ouest) : Absence totale de réservoir utile (Net Pay 0 m). Pincement et colmatage diagénétique."
    ]
    add_bullet_points(s5, pts5, Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.0), size=15)
    try_add_image(s5, os.path.join("outputs", "plots", "Synthese_Performances_Puits.png"), Inches(1.5), Inches(3.6), Inches(10.333), Inches(3.5))

    # -------------------------------------------------------------------------
    # DIAPO 6 : FOCUS TECHNIQUE PUITS GT-02
    # -------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_title(s6, "5. Le Puits Producteur GT-02 : Zone de Faille Majeure")
    pts6 = [
        "• Cible structurale : Couloir de faille sous-vertical au sein des calcaires profonds du Carbonifère.",
        "• Épaisseur totale forée : 999 mètres. Net-to-Gross confiné de 15.02% (150 m utiles).",
        "• Comportement hydrodynamique : Rôle de drain vertical ascendant ramenant des fluides hydrothermaux profonds.",
        "• Température de production : 53.74°C.",
        "• Énergie disponible en place (USGS) : 5.28 PétaJoules (PJ) localisés dans le réseau de fractures."
    ]
    add_bullet_points(s6, pts6, Inches(0.5), Inches(1.8), Inches(6.5), Inches(5.0), size=16)
    try_add_image(s6, os.path.join("outputs", "plots", "Well_GT02_Tracks.png"), Inches(7.2), Inches(1.8), Inches(5.6), Inches(4.8))

    # -------------------------------------------------------------------------
    # DIAPO 7 : CORRÉLATION STRUCTURALE INTER-PUITS (COUPE 2D)
    # -------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_title(s7, "6. Modèle Géométrique et Compartimentage du Bassin")
    try_add_image(s7, os.path.join("output", "coupe_lithostratigraphique_2D.png"), Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))

    # -------------------------------------------------------------------------
    # DIAPO 8 : FONDEMENTS PHYSIQUES & ÉQUATIONS RÉSERVOIR
    # -------------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_slide_layout)
    add_title(s8, "7. Justification Physique et Mathématique")
    pts8 = [
        "1. Modèle Volumétrique du Contenu Thermique (Méthode USGS) :",
        " Q_HIP = A * h * [((1 - Phi) * Rhome * Cme) + (Phi * Rhof * Cf)] * (Tr - Tref)",
        " - Certifie le volume de chaleur stockée en place sans spéculation empirique.",
        " - Intègre la matrice rocheuse (2650 kg/m³, 850 J/kg/°C) et l'eau (4184 J/kg/°C). Température limite de rejet = 40°C.",
        "",
        "2. Équation de Dimensionnement Hydraulique Exigé :",
        " Q_production = P_thermal / [Rhof * Cf * (Tr - Tref)]",
        " - Traduit la puissance de l'échangeur de surface en besoin de débit dynamique réel.",
        " - Application numérique validée pour GT-02 : Débit critique de 55.18 m³/h (15.33 L/s)."
    ]
    add_bullet_points(s8, pts8, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), size=15)

    # -------------------------------------------------------------------------
    # DIAPO 9 : RECOMMANDATIONS STRATÉGIQUES ET PLAN D'INGÉNIERIE
    # -------------------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_slide_layout)
    add_title(s9, "8. Plan d'Action d'Ingénierie de Production")
    pts9 = [
        "• Complétion du puits GT-02 : Recommandation d'une architecture Open-Hole ou Premium Slotted Liners pour préserver le réseau de fractures naturelles sans colmatage.",
        "• Spécifications de l'équipement de levage (ESP) : Pompe immergée calibrée pour un débit nominal de 50 à 65 m³/h, positionnée à ~350 m de profondeur.",
        "• Métallurgie stricte anti-corrosion : Utilisation d'aciers super-duplex ou alliage de Nickel (Inconel 625) en raison de la présence de CO2 dissous.",
        "• Reconversion de l'actif GT-01 : Reconversion en unité ATES (Stockage Thermique en Aquifère) pour capter et valoriser les excédents de chaleur industriels de surface en période estivale."
    ]
    add_bullet_points(s9, pts9, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0), size=16)

    # -------------------------------------------------------------------------
    # DIAPO 10 : SYNTHÈSE INVESTISSEURS & ATTRIBUTION DU RISQUE
    # -------------------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_slide_layout)
    add_title(s10, "9. Conclusion : Facteurs de Succès et Maîtrise des Risques")
    pts10 = [
        "• Risque de sous-sol drastiquement réduit : Validation déterministe via le pipeline de données réelles MySQL.",
        "• Ressource commerciale confirmée : 5.28 PJ d'énergie thermique en place identifiés sur le segment GT-02.",
        "• Viabilité dynamique démontrée : Débit cible de 55.18 m³/h en adéquation complète avec les propriétés de la faille.",
        "• Optimisation de l'investissement (CAPEX) : Double valeur d'usage du site (Production sur GT-02 et Stockage sur GT-01).",
        "• Cycle d'exploitation sécurisé : Horizon opérationnel de 30 ans sur le doublet hydrothermique."
    ]
    add_bullet_points(s10, pts10, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0), size=16)

    output_path = "Presentation_Potentiel_Geothermie_Bloc_UK.pptx"
    prs.save(output_path)
    print(f"\n🎉 Présentation PowerPoint générée avec succès : {output_path}")

if __name__ == "__main__":
    create_presentation()